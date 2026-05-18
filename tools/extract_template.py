"""원본 57장 PPTX에서 카탈로그 대상 10장만 남긴 마스터 템플릿을 생성한다.

추출 후 슬라이드 순서(1-based)와 원본 슬라이드 번호:
  1=원본20 cover     2=원본21 index     3=원본22 divider
  4=원본27 content-1col  5=원본28 content-2col
  6=원본29 content-3col  7=원본30 content-4col
  8=원본38 chart     9=원본46 table     10=원본57 closing
"""
import sys
from pptx import Presentation

SRC = "docs/reference/original-template.pptx"
DST = "skills/generate-deepnoid-ppt/assets/deepnoid-template.pptx"
KEEP_0BASED = [19, 20, 21, 26, 27, 28, 29, 37, 45, 56]


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides._sldIdLst)
    if len(slides) != 57:
        sys.exit(f"원본 슬라이드 수가 57이 아님: {len(slides)}")
    keep = set(KEEP_0BASED)
    for i, sld in enumerate(slides):
        if i not in keep:
            prs.slides._sldIdLst.remove(sld)
    prs.save(DST)
    out = Presentation(DST)
    n = len(out.slides._sldIdLst)
    if n != 10:
        sys.exit(f"추출 결과가 10장이 아님: {n}")
    print(f"추출 완료: {DST} ({n}장)")


if __name__ == "__main__":
    main()

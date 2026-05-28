"""우상단 헤더(슬로건 + 로고 + 페이지번호) — design.md §4 내지 전용."""
from pathlib import Path
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

SLOGAN_TEXT = "Through AI, We make your life wider, bolder and clearer"

# design.md §4 좌표 (cm)
SLOGAN_POS = (17.86, 0.57)
SLOGAN_SIZE = (10.42, 0.3)
LOGO_POS = (28.65, 0.54)
LOGO_SIZE = (2.71, 0.42)
PAGENUM_POS = (31.36, 0.54)
PAGENUM_SIZE = (1.19, 0.42)

MUTED = RGBColor(0x6E, 0x6E, 0x73)


def add_inner_header(slide, page_num: int, logo_path: Path) -> None:
    """슬라이드에 우상단 헤더 3요소를 추가.

    표지·간지·아웃트로에는 사용하지 않는다(내지 전용).
    """
    _add_slogan(slide)
    if Path(logo_path).exists():
        slide.shapes.add_picture(
            str(logo_path),
            left=Cm(LOGO_POS[0]), top=Cm(LOGO_POS[1]),
            width=Cm(LOGO_SIZE[0]), height=Cm(LOGO_SIZE[1]),
        )
    _add_page_number(slide, page_num)


def _add_slogan(slide) -> None:
    tb = slide.shapes.add_textbox(
        left=Cm(SLOGAN_POS[0]), top=Cm(SLOGAN_POS[1]),
        width=Cm(SLOGAN_SIZE[0]), height=Cm(SLOGAN_SIZE[1]),
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = SLOGAN_TEXT
    run.font.name = "Pretendard"
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED


def _add_page_number(slide, page_num: int) -> None:
    tb = slide.shapes.add_textbox(
        left=Cm(PAGENUM_POS[0]), top=Cm(PAGENUM_POS[1]),
        width=Cm(PAGENUM_SIZE[0]), height=Cm(PAGENUM_SIZE[1]),
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page_num)
    run.font.name = "Pretendard"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED

"""DEEPNOID 슬라이드 조립 — 표지·간지·아웃트로 + outline 디스패처."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from tools.deepnoid_builder.headers import add_inner_header
from tools.deepnoid_builder.patterns import (
    add_eyebrow_and_title, add_card_grid, add_comparison, add_step_flow, add_kpi_cards,
    BLUE, GREEN, BLACK, MUTED, _fmt_text,
)

SLIDE_W_CM = 33.87
SLIDE_H_CM = 19.05
SLOGAN = "Through AI, We make your life wider, bolder and clearer"


def new_presentation() -> Presentation:
    p = Presentation()
    p.slide_width = Cm(SLIDE_W_CM)
    p.slide_height = Cm(SLIDE_H_CM)
    return p


def _add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# ================== 표지 ==================
def add_cover(prs, title: str, subtitle: str, logo_path: Path) -> None:
    slide = _add_blank_slide(prs)
    # 제목 (대형 BLUE)
    tb = slide.shapes.add_textbox(Cm(1.91), Cm(6.0), Cm(30.0), Cm(4.0))
    _fmt_text(tb, title, size_pt=32, bold=True, color=BLUE)
    # 부제 (BLACK)
    tb = slide.shapes.add_textbox(Cm(1.91), Cm(10.5), Cm(30.0), Cm(1.0))
    _fmt_text(tb, subtitle, size_pt=14, bold=False, color=BLACK)
    # 로고
    if Path(logo_path).exists():
        slide.shapes.add_picture(str(logo_path),
                                 Cm(1.91), Cm(14.29), Cm(3.06), Cm(1.35))


# ================== 간지 ==================
def add_divider(prs, title: str, chapter: str = "", number: int | None = None) -> None:
    slide = _add_blank_slide(prs)
    # Chapter 라벨 (있을 때)
    if chapter:
        tb = slide.shapes.add_textbox(Cm(1.91), Cm(1.91), Cm(10.0), Cm(0.8))
        _fmt_text(tb, chapter, size_pt=20, bold=False, color=BLUE)
    # Main Title (36pt Bold)
    tb = slide.shapes.add_textbox(Cm(2.41), Cm(4.0), Cm(28.0), Cm(4.0))
    _fmt_text(tb, title, size_pt=36, bold=True, color=BLACK)
    # 글리프 (180pt GREEN Bold)
    if number is not None:
        tb = slide.shapes.add_textbox(Cm(2.3), Cm(12.5),
                                      Cm(8.0) if number >= 10 else Cm(3.7), Cm(3.33))
        _fmt_text(tb, str(number), size_pt=180, bold=True, color=GREEN)


# ================== 아웃트로 ==================
def add_outro(prs, logo_path: Path) -> None:
    slide = _add_blank_slide(prs)
    # 슬로건 32pt
    tb = slide.shapes.add_textbox(Cm(1.91), Cm(4.40), Cm(20.93), Cm(2.87))
    _fmt_text(tb, SLOGAN, size_pt=32, bold=True, color=BLACK)
    # 좌하단 워드마크
    if Path(logo_path).exists():
        slide.shapes.add_picture(str(logo_path),
                                 Cm(1.91), Cm(14.29), Cm(3.06), Cm(1.35))
    # 푸터
    tb = slide.shapes.add_textbox(Cm(1.91), Cm(16.08), Cm(17.03), Cm(1.45))
    _fmt_text(tb, "DEEPNOID Inc.  © 2026 All rights reserved.",
              size_pt=9, bold=False, color=MUTED)


LOGO_PATH = Path("skills/generate-deepnoid-ppt/assets/deepnoid_logo.png")


def build(outline: dict, logo_path: Path = LOGO_PATH) -> Presentation:
    """outline dict → Presentation 객체. 슬라이드 타입별로 디스패치한다."""
    prs = new_presentation()
    inner_page_num = 0  # 내지는 1부터 매김
    for spec in outline.get("slides", []):
        t = spec.get("type")
        if t == "cover":
            add_cover(prs, title=spec["title"], subtitle=spec.get("subtitle", ""),
                      logo_path=logo_path)
        elif t == "divider":
            add_divider(prs, title=spec["title"],
                        chapter=spec.get("chapter", ""),
                        number=spec.get("number"))
        elif t == "outro":
            add_outro(prs, logo_path=logo_path)
        elif t in ("card-grid-2", "card-grid-3", "card-grid-4", "card-grid-7"):
            inner_page_num += 1
            slide = _add_blank_slide(prs)
            add_inner_header(slide, page_num=inner_page_num, logo_path=logo_path)
            add_eyebrow_and_title(slide, spec.get("eyebrow", ""), spec.get("title", ""))
            n = int(t.split("-")[-1])
            add_card_grid(slide, n=n, cards=spec.get("cards", []),
                          note=spec.get("note", ""))
        elif t == "comparison":
            inner_page_num += 1
            slide = _add_blank_slide(prs)
            add_inner_header(slide, page_num=inner_page_num, logo_path=logo_path)
            add_eyebrow_and_title(slide, spec.get("eyebrow", ""), spec.get("title", ""))
            add_comparison(slide, asis=spec["asis"], tobe=spec["tobe"],
                           caption=spec.get("caption", ""))
        elif t == "step":
            inner_page_num += 1
            slide = _add_blank_slide(prs)
            add_inner_header(slide, page_num=inner_page_num, logo_path=logo_path)
            add_eyebrow_and_title(slide, spec.get("eyebrow", ""), spec.get("title", ""))
            add_step_flow(slide, steps=spec.get("steps", []),
                          footer=spec.get("footer", ""))
        elif t == "kpi":
            inner_page_num += 1
            slide = _add_blank_slide(prs)
            add_inner_header(slide, page_num=inner_page_num, logo_path=logo_path)
            add_eyebrow_and_title(slide, spec.get("eyebrow", ""), spec.get("title", ""))
            add_kpi_cards(slide, cards=spec.get("cards", []))
        else:
            raise ValueError(f"알 수 없는 슬라이드 type: {t}")
    return prs

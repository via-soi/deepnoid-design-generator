"""§6 내지 본문 패턴 — design.md §6.0~6.5 좌표·크기를 1:1 매핑.

함수마다 슬라이드 객체와 인자 dict 를 받아 도형을 추가한다. 외부 상태 의존 없음.
"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt

# --- design.md §2 색 ---
BLUE = RGBColor(0x00, 0x66, 0xFF)
GREEN = RGBColor(0x34, 0xD0, 0xB3)
NAVY = RGBColor(0x00, 0x00, 0x57)
BLACK = RGBColor(0x00, 0x00, 0x00)
MUTED = RGBColor(0x6E, 0x6E, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xCC, 0xE0, 0xFF)
LIGHT_GREEN = RGBColor(0xD6, 0xF6, 0xF0)
SURFACE = RGBColor(0xF2, 0xF2, 0xF2)
HAIRLINE = RGBColor(0xD1, 0xD1, 0xD6)

# --- 본문 영역 기준 ---
BODY_X = 1.33   # cm
BODY_Y = 6.5    # cm

# --- §6.0 골격 좌표 ---
EYEBROW_POS = (1.33, 1.4)
EYEBROW_SIZE = (20.0, 0.6)
PAGE_TITLE_POS = (1.33, 2.3)
PAGE_TITLE_SIZE = (28.0, 2.4)


# ================== §6.0 골격 ==================
def add_eyebrow_and_title(slide, eyebrow: str, title: str) -> None:
    """내지 상단 골격: eyebrow(11pt BLUE) + 페이지 제목(26pt Bold)."""
    # eyebrow
    tb = slide.shapes.add_textbox(
        Cm(EYEBROW_POS[0]), Cm(EYEBROW_POS[1]),
        Cm(EYEBROW_SIZE[0]), Cm(EYEBROW_SIZE[1]),
    )
    _fmt_text(tb, eyebrow, size_pt=11, bold=True, color=BLUE)
    # title
    tb = slide.shapes.add_textbox(
        Cm(PAGE_TITLE_POS[0]), Cm(PAGE_TITLE_POS[1]),
        Cm(PAGE_TITLE_SIZE[0]), Cm(PAGE_TITLE_SIZE[1]),
    )
    _fmt_text(tb, title, size_pt=26, bold=True, color=BLACK)


# ================== §6.1 카드 ==================
def add_card(slide, x_cm: float, y_cm: float, w_cm: float, h_cm: float,
             header: str, body: str, accent: str = "gray") -> None:
    """둥근 사각형 카드 1장. accent: 'gray' | 'blue' | 'green' | 'blue-solid' | 'green-solid'.

    'gray'(기본) — SURFACE 채움 + HAIRLINE 1pt 테두리
    'blue'/'green' — LIGHT_BLUE/LIGHT_GREEN 채움 + 메인색 1pt 테두리
    'blue-solid'/'green-solid' — 메인색 솔리드 채움(§6.5 KPI 전용)
    """
    fill, line, text_color = _accent_to_colors(accent)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm),
    )
    rect.adjustments[0] = 0.05  # ≈6px
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(1)
    rect.shadow.inherit = False
    # header
    inner_left = x_cm + 0.5
    inner_top = y_cm + 0.95
    inner_w = w_cm - 1.0
    tb = slide.shapes.add_textbox(Cm(inner_left), Cm(inner_top), Cm(inner_w), Cm(1.0))
    _fmt_text(tb, header, size_pt=14, bold=True, color=text_color)
    # body
    tb = slide.shapes.add_textbox(Cm(inner_left), Cm(inner_top + 1.2),
                                  Cm(inner_w), Cm(h_cm - 2.6))
    _fmt_text(tb, body, size_pt=14, bold=False, color=text_color)


def _accent_to_colors(accent: str):
    a = accent.lower()
    if a == "gray":
        return SURFACE, HAIRLINE, BLACK
    if a == "blue":
        return LIGHT_BLUE, BLUE, BLACK
    if a == "green":
        return LIGHT_GREEN, GREEN, BLACK
    if a == "blue-solid":
        return BLUE, None, WHITE
    if a == "green-solid":
        return GREEN, None, WHITE
    raise ValueError(f"알 수 없는 accent: {accent}")


# ================== §6.2 카드 그리드 ==================
GRID_SPECS = {
    # n: (card_w, card_h, gap_x)
    2: (15.20, 9.0, 0.50),
    3: (10.07, 9.0, 0.50),
    4: (7.46, 9.0, 0.45),
}
GRID_7 = (7.46, 5.3, 0.45, 0.40)  # w, h, gap_x, gap_y — 4×2


def add_card_grid(slide, n: int, cards: list, note: str = "") -> None:
    """§6.2 카드 그리드. n ∈ {2,3,4,7}.

    cards: [{"header": str, "body": str, "accent": "gray"|"blue"|"green"}, ...]
    note (n=7 한정): 8번째 안내 카드의 본문. 헤더 없음.
    """
    if n == 7:
        return _grid_7(slide, cards, note)
    if n not in GRID_SPECS:
        raise ValueError(f"지원하지 않는 카드 수: {n}")
    w, h, gap = GRID_SPECS[n]
    if len(cards) < n:
        cards = cards + [{"header": "", "body": "", "accent": "gray"}] * (n - len(cards))
    for i in range(n):
        x = BODY_X + i * (w + gap)
        c = cards[i]
        add_card(slide, x_cm=x, y_cm=BODY_Y, w_cm=w, h_cm=h,
                 header=c.get("header", ""), body=c.get("body", ""),
                 accent=c.get("accent", "gray"))


def _grid_7(slide, cards: list, note: str) -> None:
    """4×2 그리드: 7장 + 1 안내 카드."""
    w, h, gx, gy = GRID_7
    if len(cards) < 7:
        cards = cards + [{"header": "", "body": "", "accent": "gray"}] * (7 - len(cards))
    for i in range(7):
        col = i % 4
        row = i // 4
        x = BODY_X + col * (w + gx)
        y = BODY_Y + row * (h + gy)
        c = cards[i]
        add_card(slide, x_cm=x, y_cm=y, w_cm=w, h_cm=h,
                 header=c.get("header", ""), body=c.get("body", ""),
                 accent=c.get("accent", "gray"))
    # 8번째: 안내 카드 (LIGHT_GREEN 배경)
    col = 3
    row = 1
    x = BODY_X + col * (w + gx)
    y = BODY_Y + row * (h + gy)
    add_card(slide, x_cm=x, y_cm=y, w_cm=w, h_cm=h,
             header="", body=note, accent="green")


# ================== §6.3 비교형 ==================
COMPARE_LEFT = (1.33, 6.5, 15.2, 8.5)    # x, y, w, h
COMPARE_RIGHT = (17.34, 6.5, 15.2, 8.5)
CAPTION_Y = 16.0


def add_comparison(slide, asis: dict, tobe: dict, caption: str = "") -> None:
    """§6.3 AS-IS / TO-BE 비교 카드 2장 + 선택 캡션.

    asis/tobe: {"label": str, "header": str, "bullets": list[str]}
    """
    # 좌측 (AS-IS)
    _compare_panel(slide, *COMPARE_LEFT,
                   fill=LIGHT_BLUE, line=None,
                   label=asis["label"], label_color=NAVY,
                   header=asis["header"], bullets=asis["bullets"],
                   highlight_last=False)
    # 우측 (TO-BE)
    _compare_panel(slide, *COMPARE_RIGHT,
                   fill=WHITE, line=(GREEN, 1.5),
                   label=tobe["label"], label_color=GREEN,
                   header=tobe["header"], bullets=tobe["bullets"],
                   highlight_last=True)
    if caption:
        tb = slide.shapes.add_textbox(Cm(BODY_X), Cm(CAPTION_Y), Cm(31.21), Cm(0.8))
        from pptx.enum.text import PP_ALIGN
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        run = p.add_run()
        run.text = caption
        run.font.name = "Pretendard"
        run.font.size = Pt(14)
        run.font.color.rgb = MUTED


def _compare_panel(slide, x, y, w, h, fill, line,
                   label: str, label_color, header: str, bullets: list,
                   highlight_last: bool) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Cm(x), Cm(y), Cm(w), Cm(h))
    rect.adjustments[0] = 0.05
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        color, width_pt = line
        rect.line.color.rgb = color
        rect.line.width = Pt(width_pt)
    rect.shadow.inherit = False
    # 라벨 (10pt Bold)
    tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(y + 0.5), Cm(w - 1.0), Cm(0.6))
    _fmt_text(tb, label, size_pt=10, bold=True, color=label_color)
    # 헤더 (14pt Bold)
    tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(y + 1.3), Cm(w - 1.0), Cm(1.0))
    _fmt_text(tb, header, size_pt=14, bold=True, color=BLACK)
    # 본문 bullets (14pt Regular)
    tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(y + 2.6),
                                  Cm(w - 1.0), Cm(h - 3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.1)
    tf.margin_top = tf.margin_bottom = Cm(0.1)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.name = "Pretendard"
        run.font.size = Pt(14)
        is_last_highlight = highlight_last and i == len(bullets) - 1
        run.font.bold = is_last_highlight
        run.font.color.rgb = BLUE if is_last_highlight else BLACK


# ================== §6.4 스텝형 ==================
from pptx.enum.shapes import MSO_CONNECTOR

STEP_W = 7.46
STEP_H = 9.0
STEP_GAP = 0.45


def add_step_flow(slide, steps: list, footer: str = "") -> None:
    """§6.4 스텝형 (3~5개 카드 + 연결 화살표).

    steps: [{"label": "STEP 1", "header": "...", "body": "..."}, ...]
    """
    n = len(steps)
    if not 3 <= n <= 5:
        raise ValueError(f"스텝 수는 3~5: 받은 값 {n}")
    for i, s in enumerate(steps):
        x = BODY_X + i * (STEP_W + STEP_GAP)
        color = BLUE if i % 2 == 0 else GREEN
        _step_card(slide, x, BODY_Y, STEP_W, STEP_H, color,
                   s.get("label", ""), s.get("header", ""), s.get("body", ""))
        # 화살표 (마지막 카드 뒤에는 없음)
        if i < n - 1:
            x_from = x + STEP_W + 0.05
            x_to = x + STEP_W + STEP_GAP - 0.05
            y_mid = BODY_Y + STEP_H / 2
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Cm(x_from), Cm(y_mid), Cm(x_to), Cm(y_mid),
            )
            conn.line.color.rgb = BLUE
            conn.line.width = Pt(0.75)
            # 화살촉 (tailEnd)
            from lxml import etree
            ln = conn.get_or_add_ln()
            etree.SubElement(
                ln,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd',
                {'type': 'triangle', 'w': 'med', 'len': 'med'},
            )
    if footer:
        tb = slide.shapes.add_textbox(Cm(BODY_X), Cm(CAPTION_Y), Cm(31.21), Cm(0.8))
        from pptx.enum.text import PP_ALIGN
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ""
        run = p.add_run()
        run.text = footer
        run.font.name = "Pretendard"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = BLUE


def _step_card(slide, x, y, w, h, color: RGBColor,
               label: str, header: str, body: str) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Cm(x), Cm(y), Cm(w), Cm(h))
    rect.adjustments[0] = 0.05
    rect.fill.solid()
    rect.fill.fore_color.rgb = SURFACE
    rect.line.color.rgb = color
    rect.line.width = Pt(1)
    rect.shadow.inherit = False
    # STEP 라벨
    tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(y + 0.5), Cm(w - 1.0), Cm(0.6))
    _fmt_text(tb, label, size_pt=10, bold=True, color=color)
    # 헤더
    tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(y + 1.4), Cm(w - 1.0), Cm(1.6))
    _fmt_text(tb, header, size_pt=14, bold=True, color=BLACK)
    # 본문
    tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(y + 3.4),
                                  Cm(w - 1.0), Cm(h - 4.0))
    _fmt_text(tb, body, size_pt=14, bold=False, color=BLACK)


# ================== §6.5 KPI·메시지형 ==================
def add_kpi_cards(slide, cards: list) -> None:
    """§6.5 KPI/메시지 카드 — 메인색 솔리드 채움, 흰 텍스트.

    cards: [{"number": "01", "header": "...", "body": "..."}, ...]
    """
    n = len(cards)
    if not 3 <= n <= 4:
        raise ValueError(f"KPI 카드 수는 3~4: 받은 값 {n}")
    w, h, gap = STEP_W, STEP_H, STEP_GAP   # 4단 그리드와 동일
    for i, c in enumerate(cards):
        x = BODY_X + i * (w + gap)
        fill = BLUE if i % 2 == 0 else GREEN
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Cm(x), Cm(BODY_Y), Cm(w), Cm(h))
        rect.adjustments[0] = 0.05
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.fill.background()
        rect.shadow.inherit = False
        # 큰 숫자
        tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(BODY_Y + 0.5), Cm(w - 1.0), Cm(2.5))
        _fmt_text(tb, c.get("number", ""), size_pt=42, bold=True, color=WHITE)
        # 헤더
        tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(BODY_Y + 3.5), Cm(w - 1.0), Cm(1.6))
        _fmt_text(tb, c.get("header", ""), size_pt=14, bold=True, color=WHITE)
        # 본문
        tb = slide.shapes.add_textbox(Cm(x + 0.5), Cm(BODY_Y + 5.5),
                                      Cm(w - 1.0), Cm(h - 6.0))
        _fmt_text(tb, c.get("body", ""), size_pt=14, bold=False, color=WHITE)


# ================== 공용 텍스트 헬퍼 ==================
def _fmt_text(shape, text: str, size_pt: float, bold: bool, color: RGBColor) -> None:
    """텍스트박스/도형에 단일 단락·단일 run 텍스트를 채운다."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    p = tf.paragraphs[0]
    p.text = ""  # 기존 run 제거
    run = p.add_run()
    run.text = text
    run.font.name = "Pretendard"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color

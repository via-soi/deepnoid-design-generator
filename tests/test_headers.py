from pathlib import Path
from pptx.util import Cm
from tools.deepnoid_builder.headers import add_inner_header

LOGO = Path("skills/generate-deepnoid-ppt/assets/deepnoid_logo.png")


def test_adds_three_shapes(slide):
    add_inner_header(slide, page_num=3, logo_path=LOGO)
    # 슬로건 txt, 로고 pic, 페이지번호 txt = 3 도형
    assert len(slide.shapes) == 3


def test_slogan_position_and_text(slide):
    add_inner_header(slide, page_num=1, logo_path=LOGO)
    slogan = [s for s in slide.shapes if s.has_text_frame and "Through AI" in s.text_frame.text][0]
    assert abs(slogan.left - Cm(17.86)) < Cm(0.05)
    assert abs(slogan.top - Cm(0.57)) < Cm(0.05)


def test_page_number(slide):
    add_inner_header(slide, page_num=7, logo_path=LOGO)
    nums = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip() == "7"]
    assert len(nums) == 1

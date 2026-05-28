from pathlib import Path
from pptx import Presentation
from tools.deepnoid_builder.builder import add_cover, add_divider, add_outro

LOGO = Path("skills/generate-deepnoid-ppt/assets/deepnoid_logo.png")


def test_cover_has_title_and_subtitle(prs):
    add_cover(prs, title="제목", subtitle="부제 · 발표자", logo_path=LOGO)
    slide = prs.slides[0]
    txts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "제목" in txts
    assert "부제" in txts


def test_divider_has_title_and_number(prs):
    add_divider(prs, title="섹션 제목", chapter="Chapter 01", number=1)
    slide = prs.slides[0]
    txts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "섹션 제목" in txts
    assert "Chapter 01" in txts
    assert "1" in txts


def test_outro_has_slogan(prs):
    add_outro(prs, logo_path=LOGO)
    slide = prs.slides[0]
    txts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Through AI" in txts


from tools.deepnoid_builder.builder import build


def test_build_full_outline():
    outline = {
        "deck": {"title": "T", "subtitle": "S"},
        "slides": [
            {"type": "cover", "title": "표지", "subtitle": "부제"},
            {"type": "divider", "chapter": "Chapter 01", "title": "섹션", "number": 1},
            {"type": "card-grid-3",
             "eyebrow": "Intro / X",
             "title": "내지 제목",
             "cards": [
                 {"header": "A", "body": "a"},
                 {"header": "B", "body": "b", "accent": "blue"},
                 {"header": "C", "body": "c"},
             ]},
            {"type": "outro"},
        ],
    }
    prs = build(outline)
    assert len(prs.slides) == 4

"""공용 pytest 픽스처."""
import pytest
from pptx import Presentation
from pptx.util import Cm


@pytest.fixture
def prs():
    """16:9 widescreen 빈 프레젠테이션."""
    p = Presentation()
    p.slide_width = Cm(33.87)
    p.slide_height = Cm(19.05)
    return p


@pytest.fixture
def slide(prs):
    """빈 슬라이드 1장 (blank layout)."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

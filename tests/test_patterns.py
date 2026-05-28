from pptx.util import Cm
from tools.deepnoid_builder.patterns import add_eyebrow_and_title, add_card, add_card_grid, add_comparison


def test_eyebrow_and_title_adds_two_textboxes(slide):
    add_eyebrow_and_title(slide, eyebrow="Intro / 지금 우리", title="도입은 모두에게")
    txts = [s for s in slide.shapes if s.has_text_frame]
    eyebrow = [s for s in txts if "Intro" in s.text_frame.text]
    title = [s for s in txts if "도입은" in s.text_frame.text]
    assert len(eyebrow) == 1
    assert len(title) == 1


def test_card_adds_rectangle_with_header_and_body(slide):
    add_card(slide, x_cm=1.33, y_cm=6.5, w_cm=10.0, h_cm=9.0,
             header="R&D", body="70%", accent="blue")
    # 1 사각형(카드) + 1 텍스트박스(헤더+본문) 또는 헤더/본문 각각 = 도형 ≥ 2
    assert len(slide.shapes) >= 2
    # 카드 안에 R&D 와 70% 텍스트가 있어야 함
    all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "R&D" in all_text
    assert "70%" in all_text


def test_grid_3_adds_three_cards(slide):
    cards = [
        {"header": "A", "body": "a", "accent": "gray"},
        {"header": "B", "body": "b", "accent": "blue"},
        {"header": "C", "body": "c", "accent": "gray"},
    ]
    add_card_grid(slide, n=3, cards=cards)
    rects = [s for s in slide.shapes if s.shape_type == 1]  # AUTO_SHAPE
    assert len(rects) == 3


def test_grid_4_adds_four_cards(slide):
    cards = [{"header": f"H{i}", "body": f"b{i}"} for i in range(4)]
    add_card_grid(slide, n=4, cards=cards)
    rects = [s for s in slide.shapes if s.shape_type == 1]
    assert len(rects) == 4


def test_grid_7_adds_eight_cards(slide):
    """7카드 + 1 안내 카드 = 도형 8."""
    cards = [{"header": f"H{i}", "body": f"b{i}"} for i in range(7)]
    add_card_grid(slide, n=7, cards=cards, note="설명 카드")
    rects = [s for s in slide.shapes if s.shape_type == 1]
    assert len(rects) == 8


def test_comparison_adds_two_panels(slide):
    add_comparison(slide,
                   asis={"label": "AS-IS · 지금", "header": "반복 업무", "bullets": ["a", "b", "c"]},
                   tobe={"label": "TO-BE · 이후", "header": "검토·책임", "bullets": ["d", "e", "f"]})
    rects = [s for s in slide.shapes if s.shape_type == 1]
    assert len(rects) == 2  # 좌·우 패널


def test_comparison_caption(slide):
    add_comparison(slide,
                   asis={"label": "L", "header": "Lh", "bullets": ["1", "2", "3"]},
                   tobe={"label": "R", "header": "Rh", "bullets": ["4", "5", "6"]},
                   caption="요약 캡션")
    cap = [s for s in slide.shapes if s.has_text_frame and "요약 캡션" in s.text_frame.text]
    assert len(cap) == 1


from tools.deepnoid_builder.patterns import add_step_flow


def test_step_flow_4_cards(slide):
    steps = [
        {"label": f"STEP {i+1}", "header": f"H{i+1}", "body": f"B{i+1}"}
        for i in range(4)
    ]
    add_step_flow(slide, steps=steps, footer="반복 핵심 메시지")
    rects = [s for s in slide.shapes if s.shape_type == 1]
    assert len(rects) == 4
    # 카드 사이 화살표 3개
    conns = [s for s in slide.shapes if s.shape_type == 9]  # LINE / connector
    assert len(conns) >= 3


def test_step_flow_footer(slide):
    steps = [{"label": "S1", "header": "H", "body": "B"}] * 3
    add_step_flow(slide, steps=steps, footer="요약 한 줄")
    footer = [s for s in slide.shapes if s.has_text_frame and "요약 한 줄" in s.text_frame.text]
    assert len(footer) == 1


from tools.deepnoid_builder.patterns import add_kpi_cards


def test_kpi_4_cards(slide):
    cards = [
        {"number": f"0{i+1}", "header": f"H{i+1}", "body": f"B{i+1}"}
        for i in range(4)
    ]
    add_kpi_cards(slide, cards=cards)
    rects = [s for s in slide.shapes if s.shape_type == 1]
    assert len(rects) == 4
